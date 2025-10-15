from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_traffic_management_presentation():
    # Create presentation
    prs = Presentation()
    
    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "🚦 Traffic Management System"
    subtitle.text = "AI-Powered Traffic Analysis & Monitoring\nBuilt with YOLOv8 & Streamlit\nReal-time Vehicle Detection • Pedestrian Monitoring • Traffic Optimization"
    
    # Format title
    title_font = title.text_frame.paragraphs[0].font
    title_font.size = Pt(44)
    title_font.bold = True
    title_font.color.rgb = RGBColor(44, 62, 80)  # Dark blue
    
    # Format subtitle
    subtitle_font = subtitle.text_frame.paragraphs[0].font
    subtitle_font.size = Pt(18)
    subtitle_font.color.rgb = RGBColor(52, 73, 94)
    
    # Slide 2: Project Overview
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "📋 Project Overview"
    
    text_frame = content.text_frame
    text_frame.clear()
    
    p = text_frame.paragraphs[0]
    p.text = "🎯 Objective: Develop an intelligent traffic management system using computer vision and AI"
    p.level = 0
    
    p = text_frame.add_paragraph()
    p.text = "🔧 Technology: YOLOv8 object detection with Streamlit web interface"
    p.level = 0
    
    p = text_frame.add_paragraph()
    p.text = "📊 Capabilities: Real-time analysis, adaptive control, violation detection, and comprehensive reporting"
    p.level = 0
    
    p = text_frame.add_paragraph()
    p.text = "🌟 Impact: Improved traffic flow, enhanced safety, and data-driven traffic management"
    p.level = 0
    
    # Slide 3: Key Features
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "🚀 Key Features"
    
    # Add text boxes for features
    left = Inches(0.5)
    top = Inches(2)
    width = Inches(4.5)
    height = Inches(2.5)
    
    # Feature 1: Adaptive Traffic Signals
    textbox1 = slide.shapes.add_textbox(left, top, width, height)
    tf1 = textbox1.text_frame
    tf1.text = "🚥 Adaptive Traffic Signals\n• Real-time vehicle density detection\n• Dynamic signal timing optimization\n• Lane-by-lane traffic analysis\n• Historical trend visualization"
    
    # Feature 2: Violation Detection
    textbox2 = slide.shapes.add_textbox(left + Inches(5), top, width, height)
    tf2 = textbox2.text_frame
    tf2.text = "⚠️ Violation Detection\n• Red light violation monitoring\n• Automated violation logging\n• Timestamp and vehicle tracking\n• Compliance reporting"
    
    # Feature 3: Vehicle Classification
    textbox3 = slide.shapes.add_textbox(left, top + Inches(3), width, height)
    tf3 = textbox3.text_frame
    tf3.text = "🚗 Vehicle Classification\n• Multi-class vehicle detection\n• Real-time counting & categorization\n• Statistical distribution analysis\n• Traffic composition insights"
    
    # Feature 4: Pedestrian Monitoring
    textbox4 = slide.shapes.add_textbox(left + Inches(5), top + Inches(3), width, height)
    tf4 = textbox4.text_frame
    tf4.text = "🚶 Pedestrian Monitoring\n• Crosswalk zone monitoring\n• Safety alert system\n• Traffic light integration\n• Pedestrian flow analysis"
    
    # Slide 4: Technical Architecture
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "🏗️ Technical Architecture"
    
    # AI/ML Components
    textbox1 = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(4.5), Inches(3))
    tf1 = textbox1.text_frame
    tf1.text = "🧠 AI/ML Components\n• YOLOv8: State-of-the-art object detection\n• OpenCV: Computer vision processing\n• Real-time Analysis: Frame-by-frame detection\n• Multi-class Detection: Vehicles & pedestrians\n\n📊 Data Processing\n• Pandas: Data analysis & CSV generation\n• NumPy: Numerical computations\n• Plotly: Interactive visualizations"
    
    # Web Interface & Data Management
    textbox2 = slide.shapes.add_textbox(Inches(5.5), Inches(2), Inches(4.5), Inches(3))
    tf2 = textbox2.text_frame
    tf2.text = "🌐 Web Interface\n• Streamlit: Interactive web application\n• Multi-page Navigation: Feature-based organization\n• Real-time Updates: Live data visualization\n• Export Capabilities: CSV & JSON reports\n\n💾 Data Management\n• CSV Logging: Structured data storage\n• Report Generation: Automated analytics\n• Historical Tracking: Trend analysis"
    
    # Slide 5: Technology Stack
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "💻 Technology Stack"
    
    content = slide.placeholders[1]
    text_frame = content.text_frame
    text_frame.clear()
    
    p = text_frame.paragraphs[0]
    p.text = "Core Technologies:"
    p.font.bold = True
    p.font.size = Pt(18)
    
    p = text_frame.add_paragraph()
    p.text = "• Python 3.8+ • YOLOv8 (Ultralytics) • Streamlit • OpenCV"
    p.level = 0
    
    p = text_frame.add_paragraph()
    p.text = "Data & Visualization:"
    p.font.bold = True
    p.font.size = Pt(18)
    
    p = text_frame.add_paragraph()
    p.text = "• Pandas • NumPy • Plotly • Pillow"
    p.level = 0
    
    p = text_frame.add_paragraph()
    p.text = "Key Features:"
    p.font.bold = True
    p.font.size = Pt(18)
    
    features = ["GPU Acceleration Support", "Real-time Processing", "Multi-format Input Support", 
                "Interactive Web Interface", "Comprehensive Reporting", "Modular Architecture"]
    
    for feature in features:
        p = text_frame.add_paragraph()
        p.text = f"✅ {feature}"
        p.level = 0
    
    # Slide 6: System Capabilities
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "⚡ System Capabilities"
    
    # Input Methods & Analytics
    textbox1 = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(4.5), Inches(4))
    tf1 = textbox1.text_frame
    tf1.text = "🎥 Input Methods\n• Video Upload: MP4, AVI, MOV, MKV\n• Image Upload: JPG, PNG, BMP\n• Live Webcam: Real-time processing\n• Batch Processing: Multiple file support\n\n📈 Analytics Features\n• Time-series analysis\n• Statistical summaries\n• Interactive charts & graphs\n• Trend identification"
    
    # Configuration & Export
    textbox2 = slide.shapes.add_textbox(Inches(5.5), Inches(2), Inches(4.5), Inches(4))
    tf2 = textbox2.text_frame
    tf2.text = "🔧 Configuration Options\n• Adjustable confidence thresholds\n• Customizable detection zones\n• Traffic light state integration\n• Alert system configuration\n\n📤 Export & Reporting\n• CSV data exports\n• JSON summary reports\n• Real-time downloads\n• Historical data tracking"
    
    # Slide 7: Implementation Highlights
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "🛠️ Implementation Highlights"
    
    content = slide.placeholders[1]
    text_frame = content.text_frame
    text_frame.clear()
    
    highlights = [
        ("🧩 Modular Design", "Organized into separate modules for each feature with shared utilities"),
        ("🎛️ User-Friendly Interface", "Streamlit-based web interface with sidebar navigation and interactive visualizations"),
        ("📊 Real-time Processing", "Frame-by-frame video analysis with progress tracking and live statistics"),
        ("🔄 Scalable Architecture", "Designed for extensibility with plugin-style feature modules and comprehensive logging")
    ]
    
    for i, (highlight, description) in enumerate(highlights):
        p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
        p.text = f"{highlight}: {description}"
        p.level = 0
        p.font.size = Pt(14)
    
    # Slide 8: Results & Benefits
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "📊 Results & Benefits"
    
    # Performance & Traffic Optimization
    textbox1 = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(4.5), Inches(4))
    tf1 = textbox1.text_frame
    tf1.text = "🎯 Performance Metrics\n• Detection Accuracy: High-precision YOLOv8 model\n• Real-time Processing: Efficient frame analysis\n• Multi-class Recognition: Vehicles & pedestrians\n• Scalable Processing: GPU acceleration support\n\n📈 Traffic Optimization\n• Adaptive signal timing based on density\n• Lane-specific traffic analysis\n• Historical trend identification\n• Data-driven decision making"
    
    # Safety & Operational Benefits
    textbox2 = slide.shapes.add_textbox(Inches(5.5), Inches(2), Inches(4.5), Inches(4))
    tf2 = textbox2.text_frame
    tf2.text = "🛡️ Safety Improvements\n• Automated violation detection\n• Pedestrian safety monitoring\n• Real-time alert system\n• Comprehensive safety reporting\n\n💼 Operational Benefits\n• Reduced manual monitoring\n• Automated report generation\n• Cost-effective solution\n• Easy deployment & maintenance"
    
    # Slide 9: Future Enhancements
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "🚀 Future Enhancements"
    
    # AI Features & Integration
    textbox1 = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(4.5), Inches(2.5))
    tf1 = textbox1.text_frame
    tf1.text = "🤖 Advanced AI Features\n• Vehicle tracking across frames\n• Speed estimation algorithms\n• Predictive traffic modeling\n• Behavior pattern analysis\n\n🌐 Integration Capabilities\n• Smart city infrastructure\n• Traffic management systems\n• Emergency services integration\n• Mobile app connectivity"
    
    # Interface & Cloud
    textbox2 = slide.shapes.add_textbox(Inches(5.5), Inches(2), Inches(4.5), Inches(2.5))
    tf2 = textbox2.text_frame
    tf2.text = "📱 Enhanced Interface\n• Mobile-responsive design\n• Dashboard customization\n• Multi-user support\n• Role-based access control\n\n☁️ Cloud & Deployment\n• Cloud-based processing\n• Edge computing support\n• Containerized deployment\n• API development"
    
    # Slide 10: Conclusion
    slide_layout = prs.slide_layouts[0]  # Title slide layout for conclusion
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "🎯 Conclusion"
    subtitle.text = "✅ Successfully developed a comprehensive AI-powered traffic management system\n✅ Integrated cutting-edge technologies including YOLOv8 and Streamlit\n✅ Delivered multiple features for traffic optimization and safety monitoring\n✅ Created user-friendly interface with real-time analytics and reporting\n✅ Designed scalable architecture for future enhancements\n\nReady for deployment and real-world testing! 🚀"
    
    # Format conclusion text
    subtitle_font = subtitle.text_frame.paragraphs[0].font
    subtitle_font.size = Pt(16)
    subtitle_font.color.rgb = RGBColor(52, 73, 94)
    
    # Save presentation
    prs.save('Traffic_Management_System_Presentation.pptx')
    print("PowerPoint presentation created successfully!")

if __name__ == "__main__":
    create_traffic_management_presentation()
